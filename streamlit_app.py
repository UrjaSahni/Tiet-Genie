import os
import tempfile
import base64
import streamlit as st
from dotenv import load_dotenv
from langchain_community.document_loaders import (
    PyPDFLoader,
    TextLoader,
    UnstructuredWordDocumentLoader,
    UnstructuredMarkdownLoader
)
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_together import ChatTogether
from pptx import Presentation
from fpdf import FPDF
import io

# ---------------- SETUP ----------------
load_dotenv()
together_api_key = os.getenv("TOGETHER_API_KEY")
st.set_page_config(page_title="Tiet-Genie 🤖", layout="wide")


# Custom styling for better text visibility
st.markdown("""
<style>
.stChatMessageContent, .stMarkdown {
    color: #111 !important;
    font-weight: 500;
}
</style>
""", unsafe_allow_html=True)

# ---------------- SIDEBAR ----------------
with st.sidebar:
    # Check if logo exists before displaying
    if os.path.exists("TIETlogo.png"):
        st.image("TIETlogo.png", width=120)
    else:
        st.markdown("## 🏫 TIET")
    
    st.markdown("## 🤖 Tiet-Genie")
    st.markdown("How can I assist you today? 😊")
    uploaded_files = st.file_uploader(
        "📎 Upload PDFs, DOCX, PPTX, TXT, or MD",
        type=["pdf", "docx", "pptx", "txt", "md"],
        accept_multiple_files=True
    )


# ---------------- LOAD DEFAULT PDFs ----------------
@st.cache_resource(show_spinner="Loading default PDFs...")
def load_default_vectorstore():
    default_files = ["rules.pdf", "AcademicRegulations.pdf"]
    docs = []
    
    # Only load files that actually exist
    for path in default_files:
        if os.path.exists(path):
            try:
                docs.extend(PyPDFLoader(path).load())
                # st.success(f"✅ Loaded {path}")  # Hidden to avoid clutter
            except Exception as e:
                st.warning(f"⚠️ Failed to load {path}: {str(e)}")
        else:
            st.warning(f"⚠️ File not found: {path}")
    
    # If no default documents were loaded, create a minimal vectorstore
    if not docs:
        st.info("ℹ️ No default PDFs found. Upload your own files to get started.")
        # Create a dummy document to initialize the vectorstore
        from langchain.docstore.document import Document
        dummy_doc = Document(page_content="Welcome to Tiet-Genie! Please upload your documents to get started.", metadata={})
        docs = [dummy_doc]

    splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
    chunks = splitter.split_documents(docs)
    embed = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
    return FAISS.from_documents(chunks, embed)


# Initialize vector store with error handling
try:
    vector_store = load_default_vectorstore()
except Exception as e:
    st.error(f"❌ Error initializing vector store: {str(e)}")
    st.stop()


# ---------------- HANDLE USER FILE UPLOADS ----------------
def load_file_to_docs(file_path, ext):
    try:
        if ext == "pdf":
            return PyPDFLoader(file_path).load()
        elif ext == "docx":
            return UnstructuredWordDocumentLoader(file_path).load()
        elif ext == "pptx":
            prs = Presentation(file_path)
            text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
            temp_txt = file_path + ".txt"
            with open(temp_txt, "w", encoding="utf-8") as f:
                f.write(text)
            return TextLoader(temp_txt).load()
        elif ext == "txt":
            return TextLoader(file_path, encoding="utf-8").load()
        elif ext == "md":
            return UnstructuredMarkdownLoader(file_path).load()
        return []
    except Exception as e:
        st.error(f"❌ Error loading file {file_path}: {str(e)}")
        return []


if uploaded_files:
    new_docs = []
    for f in uploaded_files:
        ext = f.name.split(".")[-1].lower()
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=f".{ext}") as tmp:
                tmp.write(f.read())
                tmp_path = tmp.name
            
            loaded_docs = load_file_to_docs(tmp_path, ext)
            if loaded_docs:
                new_docs.extend(loaded_docs)
                st.success(f"✅ Successfully loaded {f.name}")
            else:
                st.warning(f"⚠️ No content extracted from {f.name}")
                
            # Clean up temporary file
            try:
                os.unlink(tmp_path)
            except:
                pass
                
        except Exception as e:
            st.error(f"❌ Error processing {f.name}: {str(e)}")

    # Only update vector store if we have new documents
    if new_docs:
        try:
            splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
            new_chunks = splitter.split_documents(new_docs)
            embed = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
            new_vs = FAISS.from_documents(new_chunks, embed)
            vector_store.merge_from(new_vs)
            st.success(f"✅ Added {len(new_docs)} documents to knowledge base")
        except Exception as e:
            st.error(f"❌ Error updating vector store: {str(e)}")


# ---------------- LLM + RETRIEVER ----------------
try:
    retriever = vector_store.as_retriever(
        search_type="mmr",
        search_kwargs={"k": 4, "fetch_k": 10, "lambda_mult": 0.5}
    )

    # Check if Together API key is available
    if not together_api_key:
        st.error("❌ TOGETHER_API_KEY not found in environment variables. Please set it in your .env file.")
        st.stop()

    llm = ChatTogether(
        model="deepseek-ai/DeepSeek-V3",
        temperature=0.2,
        together_api_key=together_api_key
    )
except Exception as e:
    st.error(f"❌ Error initializing LLM or retriever: {str(e)}")
    st.stop()


# ---------------- CHAT HISTORY ----------------
if "chat_history" not in st.session_state:
    st.session_state.chat_history = []
if "greeted" not in st.session_state:
    st.session_state.greeted = False

if not st.session_state.greeted and not st.session_state.chat_history:
    st.markdown("<h2 style='text-align:center;'>👋 Hello TIETian! How can I help you today?</h2>", unsafe_allow_html=True)


# ---------------- CHAT UI ----------------
for msg in st.session_state.chat_history:
    with st.chat_message(msg["role"]):
        st.markdown(msg["message"], unsafe_allow_html=True)

user_prompt = st.chat_input("Ask something about TIET or the lecture notes...")
if user_prompt:
    with st.chat_message("user"):
        st.write(user_prompt)
    st.session_state.chat_history.append({"role": "user", "message": user_prompt})

    with st.chat_message("assistant"):
        with st.spinner("Thinking..."):
            try:
                retrieved_docs = retriever.get_relevant_documents(user_prompt)

                context_text = "\n\n".join([
                    f"[Page {doc.metadata.get('page', '?')}] {doc.page_content.strip()}" for doc in retrieved_docs
                ])

                prompt_to_llm = f"""
You are an AI assistant for Thapar Institute. Use the following document snippets to answer the question. Be specific and provide detailed information, but DO NOT include page numbers or citations like [Page X] in your main answer. The source snippets will be provided separately.

--- DOCUMENTS ---
{context_text}

--- QUESTION ---
{user_prompt}
"""

                response_obj = llm.invoke(prompt_to_llm)
                response = response_obj.content.strip() if hasattr(response_obj, "content") else str(response_obj).strip()

                # Check if the response indicates no relevant information found
                no_info_indicators = [
                    "do not contain specific information",
                    "do not contain information",
                    "does not contain specific information",
                    "does not contain information",
                    "do not address",
                    "does not address",
                    "do not discuss",
                    "does not discuss",
                    "no information about",
                    "no specific information about",
                    "not mentioned in the documents",
                    "not found in the documents",
                    "documents do not mention",
                    "excerpts do not contain",
                    "snippets do not contain",
                    "provided documents do not"
                ]
                
                # Check if response indicates no relevant information
                response_lower = response.lower()
                show_sources = not any(indicator in response_lower for indicator in no_info_indicators)
                
                if show_sources:
                    # Show source snippets only if relevant information was found
                    source_section = "\n\n---\n\n**📄 Source Snippets:**\n"
                    for i, doc in enumerate(retrieved_docs, 1):
                        page = doc.metadata.get("page", "?")
                        snippet = doc.page_content.strip().replace("\n", " ")[:300]
                        source_section += f"- **Snippet {i} (Page {page})**: {snippet}\n"
                    final_response = f"{response}\n{source_section}"
                else:
                    # Don't show source snippets if no relevant information found
                    final_response = response

                st.markdown(final_response, unsafe_allow_html=True)
                st.session_state.chat_history.append({"role": "assistant", "message": final_response})

            except Exception as e:
                error_response = f"⚠️ Error: {str(e)}"
                st.markdown(error_response)
                st.session_state.chat_history.append({"role": "assistant", "message": error_response})


# ---------------- EXPORT CHAT HISTORY ----------------
def export_chat_history():
    chat = st.session_state.chat_history
    if not chat:
        return

    try:
        # --- PDF Export (Using built-in fonts only) ---
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)  # Use built-in Arial font
        pdf.set_auto_page_break(auto=True, margin=15)
        
        for msg in chat:
            role = "You" if msg["role"] == "user" else "Tiet-Genie"
            # Clean message text to handle special characters and remove markdown
            clean_message = msg['message'].replace('**', '').replace('*', '').replace('#', '')
            # Handle encoding issues
            clean_message = clean_message.encode('latin1', 'ignore').decode('latin1')
            pdf.multi_cell(0, 10, f"{role}:\n{clean_message}\n")

        # Output PDF to bytes - Fixed this part
        pdf_output = pdf.output(dest='S')
        if isinstance(pdf_output, str):
            pdf_bytes = pdf_output.encode('latin1')
        else:
            pdf_bytes = pdf_output
        
        pdf_buffer = io.BytesIO(pdf_bytes)

        # --- TXT Export ---
        txt_buffer = io.StringIO()
        for msg in chat:
            role = "You" if msg["role"] == "user" else "Tiet-Genie"
            # Clean markdown from text export too
            clean_txt = msg['message'].replace('**', '').replace('*', '').replace('#', '')
            txt_buffer.write(f"{role}:\n{clean_txt}\n\n")
        txt_buffer.seek(0)

        st.sidebar.markdown("### 📤 Export Chat History")
        st.sidebar.download_button(
            "⬇️ Download as .pdf",
            data=pdf_buffer,
            file_name="chat_history.pdf",
            mime="application/pdf"
        )
        st.sidebar.download_button(
            "⬇️ Download as .txt",
            data=txt_buffer.getvalue(),
            file_name="chat_history.txt",
            mime="text/plain"
        )
        
    except Exception as e:
        st.sidebar.error(f"Export error: {str(e)}")
        # Fallback: Only show TXT export if PDF fails
        try:
            txt_buffer = io.StringIO()
            for msg in chat:
                role = "You" if msg["role"] == "user" else "Tiet-Genie"
                # Clean markdown from fallback text too
                clean_txt = msg['message'].replace('**', '').replace('*', '').replace('#', '')
                txt_buffer.write(f"{role}:\n{clean_txt}\n\n")
            txt_buffer.seek(0)
            
            st.sidebar.markdown("### 📤 Export Chat History")
            st.sidebar.download_button(
                "⬇️ Download as .txt",
                data=txt_buffer.getvalue(),
                file_name="chat_history.txt",
                mime="text/plain"
            )
        except Exception as fallback_error:
            st.sidebar.error(f"Export failed: {str(fallback_error)}")


export_chat_history()
